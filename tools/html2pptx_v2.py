#!/usr/bin/env python3
"""html2pptx_v2.py - Semantic HTML-to-PowerPoint converter.

Uses native PowerPoint layout (auto-size, tables, multi-paragraph frames)
instead of absolute positioning with manual font metrics.

Core principles:
  - Cards = single text frame with multi-paragraph rich text
  - Tables = native add_table()
  - Card grids = tables with invisible borders, each cell is a card
  - Use MSO_AUTO_SIZE for content shapes
  - Use space_before/space_after for inter-element spacing
  - Conservative current_top advancement (overestimate > underestimate)
  - No hand-rolled font metrics
"""

import argparse
import math
import re
import sys
import warnings
from copy import deepcopy
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 5.625
CONTENT_LEFT = 0.8
CONTENT_WIDTH = 8.4
CONTENT_RIGHT = CONTENT_LEFT + CONTENT_WIDTH

GAP_TIGHT = 0.08
GAP_NORMAL = 0.12
GAP_SECTION = 0.25

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
MS_CYAN = RGBColor(0x50, 0xE6, 0xFF)
MS_GREEN = RGBColor(0x00, 0xCC, 0x6A)
MS_ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
MS_RED = RGBColor(0xFF, 0x45, 0x3A)
MS_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GRAY_70 = RGBColor(0xB3, 0xB3, 0xB3)
GRAY_50 = RGBColor(0x80, 0x80, 0x80)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
BORDER_GRAY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)
CODE_GREEN = RGBColor(0x4A, 0xDE, 0x80)
CODE_BLUE = RGBColor(0x60, 0xA5, 0xFA)
CODE_YELLOW = RGBColor(0xFB, 0xBF, 0x24)
CODE_GRAY = RGBColor(0x6B, 0x73, 0x80)
CODE_PURPLE = RGBColor(0xC0, 0x84, 0xFC)
CODE_STRING = RGBColor(0xFB, 0xBF, 0x24)
CODE_DEFAULT = RGBColor(0xE6, 0xE6, 0xE6)
DEFAULT_FONT = "Arial"
CODE_FONT = "Consolas"

# Bullet characters that indicate a list item
BULLET_CHARS = {"•", "-", "*", "\u2022", "\u2713", "\u2717", "✓", "✗", "→"}

# ---------------------------------------------------------------------------
# Reused helpers (exact copies from v1)
# ---------------------------------------------------------------------------


def hex_to_rgb(hex_str: str) -> Optional[RGBColor]:
    hex_str = hex_str.strip().lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join(c * 2 for c in hex_str)
    if len(hex_str) != 6:
        return None
    try:
        return RGBColor(
            int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
        )
    except (ValueError, TypeError):
        return None


def parse_color_from_class(classes: list[str]) -> Optional[RGBColor]:
    color_map = {
        "green": MS_GREEN,
        "orange": MS_ORANGE,
        "red": MS_RED,
        "ms-green": MS_GREEN,
        "ms-orange": MS_ORANGE,
        "ms-red": MS_RED,
        "ms-blue": MS_BLUE,
        "ms-cyan": MS_CYAN,
        "ms-purple": MS_PURPLE,
        "warning": MS_ORANGE,
    }
    for cls in classes:
        if cls in color_map:
            return color_map[cls]
    return None


def _replace_br_tags(element: Tag):
    for br in element.find_all("br"):
        br.replace_with("\n")


def get_text(element: Optional[Tag]) -> str:
    if element is None:
        return ""
    el_copy = deepcopy(element)
    _replace_br_tags(el_copy)
    text = el_copy.get_text(" ")
    lines = text.split("\n")
    lines = [" ".join(line.split()) for line in lines]
    return "\n".join(lines).strip()


def get_rich_text(element: Optional[Tag]) -> list[dict]:
    if element is None:
        return []
    el_copy = deepcopy(element)
    _replace_br_tags(el_copy)
    runs: list[dict] = []
    for child in el_copy.descendants:
        if isinstance(child, NavigableString):
            text = str(child)
            if not text:
                continue
            lines = text.split("\n")
            lines = [" ".join(line.split()) for line in lines]
            text = "\n".join(lines)
            if not text or text.isspace():
                continue
            bold = False
            italic = False
            color = None
            parent = child.parent
            while parent and parent.name:
                if parent.name in ("strong", "b"):
                    bold = True
                if parent.name in ("em", "i"):
                    italic = True
                if parent.name == "span":
                    cls = parent.get("class", [])
                    if "highlight" in cls:
                        color = MS_CYAN
                    elif "check" in cls:
                        color = MS_GREEN
                parent = parent.parent
            runs.append({"text": text, "bold": bold, "italic": italic, "color": color})
    merged: list[dict] = []
    for run in runs:
        if (
            merged
            and merged[-1]["bold"] == run["bold"]
            and merged[-1]["italic"] == run["italic"]
            and merged[-1]["color"] == run["color"]
        ):
            merged[-1]["text"] += run["text"]
        else:
            merged.append(run)
    if merged:
        merged[0]["text"] = merged[0]["text"].lstrip()
        merged[-1]["text"] = merged[-1]["text"].rstrip()
        merged = [r for r in merged if r["text"]]
    return merged


def extract_css_vars(soup: BeautifulSoup) -> dict[str, str]:
    result: dict[str, str] = {}
    for style_tag in soup.find_all("style"):
        css_text = style_tag.string or ""
        for m in re.finditer(r"--([\w-]+)\s*:\s*([^;]+)", css_text):
            result[m.group(1).strip()] = m.group(2).strip()
    return result


def resolve_accent_color(css_vars: dict[str, str]) -> RGBColor:
    for var_name in ("color-accent", "accent"):
        val = css_vars.get(var_name, "")
        if val.startswith("#"):
            rgb = hex_to_rgb(val)
            if rgb:
                return rgb
    return MS_BLUE


def set_slide_background(slide, color=BLACK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Table cell border helpers (XML workaround)
# ---------------------------------------------------------------------------


def _set_cell_border(cell, color_hex: str = "333333", width_emu: int = 12700):
    """Set borders on a table cell. width_emu: 12700 = 1pt."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.makeelement(
            qn(edge),
            {"w": str(width_emu), "cap": "flat", "cmpd": "sng", "algn": "ctr"},
        )
        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": color_hex})
        solidFill.append(srgbClr)
        ln.append(solidFill)
        prstDash = ln.makeelement(qn("a:prstDash"), {"val": "solid"})
        ln.append(prstDash)
        existing = tcPr.findall(qn(edge))
        for e in existing:
            tcPr.remove(e)
        tcPr.append(ln)


def _set_cell_no_border(cell):
    """Remove all borders from a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.makeelement(qn(edge), {"w": "0", "cap": "flat", "cmpd": "sng"})
        noFill = ln.makeelement(qn("a:noFill"), {})
        ln.append(noFill)
        existing = tcPr.findall(qn(edge))
        for e in existing:
            tcPr.remove(e)
        tcPr.append(ln)


def _set_cell_fill(cell, color: RGBColor):
    """Set solid fill on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing fill
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    for old in tcPr.findall(qn("a:noFill")):
        tcPr.remove(old)
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(
        qn("a:srgbClr"),
        {"val": f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"},
    )
    solidFill.append(srgbClr)
    tcPr.insert(0, solidFill)


def _set_cell_no_fill(cell):
    """Remove fill from a table cell (transparent)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    noFill = tcPr.makeelement(qn("a:noFill"), {})
    tcPr.insert(0, noFill)


# ---------------------------------------------------------------------------
# New primitive helpers
# ---------------------------------------------------------------------------


def _set_run_font(
    run,
    name: str = DEFAULT_FONT,
    size: float = 12,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WHITE,
):
    """Set font properties on a pptx run object."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _add_paragraph(
    tf,
    text: str,
    font_size: float,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WHITE,
    font_name: str = DEFAULT_FONT,
    alignment=None,
    space_before: float = 0,
    space_after: float = 0,
):
    """Add a paragraph to an existing text frame. Returns the paragraph.

    Uses the existing first paragraph if the text frame is empty,
    otherwise appends a new one.
    """
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    run = p.add_run()
    run.text = text
    _set_run_font(
        run, name=font_name, size=font_size, bold=bold, italic=italic, color=color
    )

    if alignment is not None:
        p.alignment = alignment
    if space_before > 0:
        p.space_before = Pt(space_before)
    if space_after > 0:
        p.space_after = Pt(space_after)

    return p


def _add_rich_paragraphs(
    tf,
    rich_runs: list[dict],
    font_size: float,
    default_color: RGBColor = WHITE,
    font_name: str = DEFAULT_FONT,
    alignment=None,
    space_before: float = 0,
    space_after: float = 0,
    bold: bool = False,
):
    """Add rich text runs as a single paragraph to an existing text frame."""
    if not rich_runs:
        return None

    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    for r in rich_runs:
        run = p.add_run()
        run.text = r["text"]
        _set_run_font(
            run,
            name=font_name,
            size=font_size,
            bold=r.get("bold", bold),
            italic=r.get("italic", False),
            color=r.get("color") or default_color,
        )

    if alignment is not None:
        p.alignment = alignment
    if space_before > 0:
        p.space_before = Pt(space_before)
    if space_after > 0:
        p.space_after = Pt(space_after)

    return p


def _make_text_shape(slide, left, top, width, min_height, auto_size="fit_shape"):
    """Create a text box and return (shape, text_frame).

    auto_size options:
      "fit_shape" -> TEXT_TO_FIT_SHAPE (text shrinks to fit)
      "fit_text"  -> SHAPE_TO_FIT_TEXT (shape grows to fit)
      "none"      -> NONE
    """
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(min_height)
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    if auto_size == "fit_shape":
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif auto_size == "fit_text":
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    else:
        tf.auto_size = MSO_AUTO_SIZE.NONE

    return shape, tf


def _make_filled_shape(
    slide,
    left,
    top,
    width,
    height,
    fill_color=DARK_GRAY,
    border_color=None,
    border_width=1,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    """Create a filled background shape. Returns shape."""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    # Reduce corner rounding for rounded rectangles
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shape.adjustments[0] = 0.05
        except Exception:
            pass
    return shape


def _split_bullet_lines(text: str) -> list[tuple[str, Optional[RGBColor]]]:
    """Split text into lines, detecting bullet/check/cross prefixes.

    Returns list of (line_text, color) tuples.
    """
    results = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        color = None
        if line.startswith(("✓", "\u2713")):
            color = MS_GREEN
        elif line.startswith(("✗", "\u2717")):
            color = MS_RED
        results.append((line, color))
    return results


def _make_card_frame(
    slide,
    left,
    top,
    width,
    min_height,
    title: str,
    body: str,
    title_size: float = 16,
    body_size: float = 12,
    title_color: RGBColor = WHITE,
    body_color: RGBColor = GRAY_70,
    accent_color: Optional[RGBColor] = None,
    fill_color: RGBColor = DARK_GRAY,
    border_color: Optional[RGBColor] = None,
    body_runs: Optional[list[dict]] = None,
):
    """Create a single card shape with title + body paragraphs.

    Returns the shape.
    """
    if accent_color and title_color == WHITE:
        title_color = accent_color

    shape = _make_filled_shape(
        slide,
        left,
        top,
        width,
        min_height,
        fill_color=fill_color,
        border_color=border_color,
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Title paragraph
    if title:
        _add_paragraph(
            tf, title, title_size, bold=True, color=title_color, space_after=4
        )

    # Body: check for bullet lines
    if body_runs:
        _add_rich_paragraphs(
            tf, body_runs, body_size, default_color=body_color, space_before=2
        )
    elif body:
        bullet_lines = _split_bullet_lines(body)
        if len(bullet_lines) > 1 and any(
            line.lstrip().startswith(tuple(BULLET_CHARS)) for line, _ in bullet_lines
        ):
            for line_text, line_color in bullet_lines:
                _add_paragraph(
                    tf,
                    line_text,
                    body_size,
                    color=line_color or body_color,
                    space_before=1,
                    space_after=1,
                )
        else:
            _add_paragraph(tf, body, body_size, color=body_color, space_before=2)

    return shape


# ---------------------------------------------------------------------------
# Main converter class
# ---------------------------------------------------------------------------


class HTMLToPPTXConverterV2:
    def __init__(self, html_content: str):
        self.soup = BeautifulSoup(html_content, "lxml")
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_WIDTH)
        self.prs.slide_height = Inches(SLIDE_HEIGHT)
        self.blank_layout = self.prs.slide_layouts[6]
        self.css_vars = extract_css_vars(self.soup)
        self.accent_color = resolve_accent_color(self.css_vars)
        self.color_vars = self._resolve_color_vars()
        self.warnings: list[str] = []

    def _resolve_color_vars(self) -> dict[str, RGBColor]:
        """Resolve CSS color variables like --blue, --success, --warning."""
        result: dict[str, RGBColor] = {}
        color_keywords = {
            "blue": MS_BLUE,
            "success": MS_GREEN,
            "warning": MS_ORANGE,
            "danger": MS_RED,
            "accent": self.accent_color,
            "text": WHITE,
            "muted": GRAY_70,
        }
        for name, val in self.css_vars.items():
            if val.startswith("#"):
                rgb = hex_to_rgb(val)
                if rgb:
                    result[name] = rgb
            elif name in color_keywords:
                result[name] = color_keywords[name]
        for kw, color in color_keywords.items():
            if kw not in result:
                result[kw] = color
        return result

    def _resolve_style_color(self, element: Tag) -> Optional[RGBColor]:
        """Extract color from inline style like style='color: var(--accent)'."""
        style = element.get("style", "") if element else ""
        if not style:
            return None
        m = re.search(r"color:\s*var\(--(\w[\w-]*)\)", style)
        if m:
            var_name = m.group(1)
            return self.color_vars.get(var_name) or self.accent_color
        m = re.search(r"(?:border-color|border):\s*var\(--(\w[\w-]*)\)", style)
        if m:
            var_name = m.group(1)
            return self.color_vars.get(var_name)
        return None

    # ------------------------------------------------------------------
    # Code extraction (reused from v1)
    # ------------------------------------------------------------------

    def _extract_code_runs(self, code_el: Tag) -> list[dict]:
        el_copy = deepcopy(code_el)
        _replace_br_tags(el_copy)
        color_map = {
            "code-keyword": CODE_BLUE,
            "keyword": CODE_BLUE,
            "code-string": CODE_STRING,
            "string": CODE_STRING,
            "code-comment": CODE_GRAY,
            "comment": CODE_GRAY,
            "code-type": CODE_GREEN,
            "type": CODE_GREEN,
            "code-func": CODE_YELLOW,
            "func": CODE_YELLOW,
            "code-number": CODE_PURPLE,
            "number": CODE_PURPLE,
            "layer-kernel": CODE_BLUE,
            "layer-foundation": CODE_GREEN,
            "layer-apps": CODE_PURPLE,
            "layer-modules": CODE_YELLOW,
        }
        runs: list[dict] = []
        for child in el_copy.descendants:
            if isinstance(child, NavigableString):
                text = str(child)
                if not text:
                    continue
                color = CODE_DEFAULT
                bold = False
                parent = child.parent
                while parent and parent != el_copy:
                    if parent.name == "span":
                        for cls in parent.get("class", []):
                            if cls in color_map:
                                color = color_map[cls]
                                break
                    if parent.name in ("strong", "b"):
                        bold = True
                    parent = parent.parent
                runs.append({"text": text, "color": color, "bold": bold})
        merged: list[dict] = []
        for r in runs:
            if (
                merged
                and merged[-1]["color"] == r["color"]
                and merged[-1]["bold"] == r["bold"]
            ):
                merged[-1]["text"] += r["text"]
            else:
                merged.append(r)
        return merged

    # ------------------------------------------------------------------
    # Slide extraction
    # ------------------------------------------------------------------

    def extract_slides(self) -> list[Tag]:
        slides = self.soup.find_all(["div", "section"], class_="slide")
        if not slides:
            slides = self.soup.find_all("section")
        return slides

    def is_centered(self, slide_div: Tag) -> bool:
        classes = slide_div.get("class", [])
        return "center" in classes or "title-slide" in classes

    # ------------------------------------------------------------------
    # Native table builder
    # ------------------------------------------------------------------

    def _make_table(
        self,
        slide,
        rows_data: list[list[str]],
        left: float,
        top: float,
        width: float,
        col_widths: Optional[list[float]] = None,
        header: bool = True,
        font_size: float = 11,
        accent_color: Optional[RGBColor] = None,
    ) -> tuple:
        """Create a native PowerPoint table. Returns (shape, estimated_height)."""
        if not rows_data:
            return None, 0

        accent = accent_color or self.accent_color
        n_rows = len(rows_data)
        n_cols = max(len(row) for row in rows_data) if rows_data else 1

        # Pad rows to uniform length
        for row in rows_data:
            while len(row) < n_cols:
                row.append("")

        row_height = Inches(0.35)
        table_shape = slide.shapes.add_table(
            n_rows,
            n_cols,
            Inches(left),
            Inches(top),
            Inches(width),
            row_height * n_rows,
        )
        table = table_shape.table

        # Set column widths
        if col_widths and len(col_widths) == n_cols:
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)
        else:
            col_w = width / n_cols
            for i in range(n_cols):
                table.columns[i].width = Inches(col_w)

        # Populate cells
        for r_idx, row in enumerate(rows_data):
            is_header_row = header and r_idx == 0
            for c_idx, cell_text in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = ""
                tf = cell.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.08)
                tf.margin_right = Inches(0.08)
                tf.margin_top = Inches(0.04)
                tf.margin_bottom = Inches(0.04)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Determine text color
                text_color = GRAY_70
                text_bold = False
                if is_header_row:
                    text_color = accent
                    text_bold = True
                elif c_idx == 0:
                    text_color = WHITE
                    text_bold = True
                # Check/cross coloring
                stripped = cell_text.strip()
                if stripped in ("\u2713", "\u2714", "\u2714\ufe0e", "Yes"):
                    text_color = MS_GREEN
                elif stripped in ("\u2717", "\u2718", "\u2715", "No"):
                    text_color = MS_RED
                elif stripped in ("~", "Partial"):
                    text_color = MS_ORANGE

                _add_paragraph(
                    tf,
                    cell_text,
                    font_size,
                    bold=text_bold,
                    color=text_color,
                )

                # Cell styling
                if is_header_row:
                    _set_cell_fill(cell, RGBColor(0x0D, 0x0D, 0x0D))
                else:
                    _set_cell_fill(
                        cell,
                        DARK_GRAY if r_idx % 2 == 1 else RGBColor(0x12, 0x12, 0x12),
                    )
                _set_cell_border(cell, "333333", 6350)  # 0.5pt

        est_height = n_rows * 0.4 + 0.2
        return table_shape, est_height

    # ------------------------------------------------------------------
    # Card grid as table
    # ------------------------------------------------------------------

    def _determine_grid_cols(self, container_classes: list[str], num_cards: int) -> int:
        """Determine column count from container CSS classes."""
        for cls in container_classes:
            if cls in ("grid-2", "halves", "split"):
                return 2
            if cls in ("grid-3", "thirds"):
                return 3
            if cls in ("grid-4", "fourths"):
                return 4
            if cls in ("grid-5",):
                return min(5, num_cards)
        return min(num_cards, 3) if num_cards > 0 else 1

    def _extract_card_content(self, card_el: Tag) -> dict:
        """Extract title, body, body_runs, and metadata from a card element."""
        classes = card_el.get("class", [])
        is_module = "module-card" in classes
        is_tool = "tool-card" in classes

        title = ""
        body = ""
        body_runs = None
        number = ""

        # Card number
        num_el = card_el.find(class_="card-number")
        if num_el:
            number = get_text(num_el)

        if is_module:
            # Module card: name, contract, purpose
            name_el = card_el.find(class_="module-name")
            contract_el = card_el.find(class_="module-contract")
            purpose_el = card_el.find(class_="module-purpose")
            title = get_text(name_el) if name_el else ""
            parts = []
            if contract_el:
                parts.append(get_text(contract_el))
            if purpose_el:
                parts.append(get_text(purpose_el))
            body = "\n".join(parts)
        elif is_tool:
            name_el = card_el.find(class_="tool-name")
            desc_el = card_el.find(class_="tool-desc")
            title = get_text(name_el) if name_el else ""
            body = get_text(desc_el) if desc_el else ""
        else:
            # Generic card
            title_el = card_el.find(
                class_=lambda c: (
                    c
                    and c in {"card-title", "card-header", "card-name"}
                )
            )
            if title_el is None:
                title_el = card_el.find(["h3", "h4"])
            if title_el:
                title = get_text(title_el)

            # Body: check for bullet lists first
            bullet_list = card_el.find("ul", class_="bullet-list") or card_el.find("ul")
            if bullet_list:
                items = bullet_list.find_all("li")
                body_lines = []
                for item in items:
                    t = get_text(item)
                    if t and not t.startswith(tuple(BULLET_CHARS)):
                        t = "• " + t
                    body_lines.append(t)
                body = "\n".join(body_lines)
            else:
                body_el = card_el.find(
                    class_=lambda c: (
                        c
                        and any(
                            x in c
                            for x in [
                                "card-body",
                                "card-desc",
                                "card-description",
                                "card-content",
                                "card-text",
                            ]
                        )
                    )
                )
                if body_el is None:
                    body_el = card_el.find("p")
                if body_el:
                    rich = get_rich_text(body_el)
                    if rich and any(
                        r.get("bold") or r.get("italic") or r.get("color") for r in rich
                    ):
                        body_runs = rich
                    else:
                        body = get_text(body_el)

        if number and title:
            title = f"{number}  {title}"
        elif number:
            title = number

        # Fallback: if no title and no body, just get all text
        if not title and not body and not body_runs:
            all_text = get_text(card_el)
            if all_text:
                lines = all_text.split("\n", 1)
                title = lines[0]
                body = lines[1] if len(lines) > 1 else ""

        # Try to resolve color
        title_color = parse_color_from_class(classes)
        style_color = self._resolve_style_color(card_el)
        if style_color:
            title_color = style_color

        return {
            "title": title,
            "body": body,
            "body_runs": body_runs,
            "title_color": title_color,
            "is_module": is_module,
        }

    def _populate_card_cell(
        self,
        cell,
        card_data: dict,
        title_size: float = 14,
        body_size: float = 11,
        accent_color: Optional[RGBColor] = None,
    ):
        """Populate a table cell as a card (title + body paragraphs)."""
        accent = accent_color or self.accent_color
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        cell.text = ""

        title = card_data.get("title", "")
        body = card_data.get("body", "")
        body_runs = card_data.get("body_runs")
        t_color = card_data.get("title_color") or accent
        is_module = card_data.get("is_module", False)

        if title:
            _add_paragraph(
                tf, title, title_size, bold=True, color=t_color, space_after=4
            )

        if is_module and body:
            # Module cards: first line is contract (monospace green), rest is purpose
            lines = body.split("\n", 1)
            _add_paragraph(
                tf,
                lines[0],
                body_size - 1,
                color=CODE_GREEN,
                font_name=CODE_FONT,
                space_before=2,
            )
            if len(lines) > 1 and lines[1].strip():
                _add_paragraph(tf, lines[1], body_size, color=GRAY_70, space_before=2)
        elif body_runs:
            _add_rich_paragraphs(
                tf,
                body_runs,
                body_size,
                default_color=GRAY_70,
                space_before=2,
            )
        elif body:
            # Check for bullet lines
            bullet_lines = _split_bullet_lines(body)
            if len(bullet_lines) > 1 and any(
                line.lstrip().startswith(tuple(BULLET_CHARS))
                for line, _ in bullet_lines
            ):
                for line_text, line_color in bullet_lines:
                    _add_paragraph(
                        tf,
                        line_text,
                        body_size,
                        color=line_color or GRAY_70,
                        space_before=1,
                        space_after=1,
                    )
            else:
                _add_paragraph(tf, body, body_size, color=GRAY_70, space_before=2)

    # ------------------------------------------------------------------
    # Element handlers
    # ------------------------------------------------------------------

    def _handle_section_number(self, slide, el, current_top, centered):
        text = get_text(el)
        if not text:
            return current_top
        top = 1.0 if centered else current_top
        shape, tf = _make_text_shape(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, 0.35, auto_size="none"
        )
        _add_paragraph(
            tf,
            text.upper(),
            14,
            bold=True,
            color=self.accent_color,
            alignment=PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT,
        )
        return top + 0.35 + GAP_TIGHT

    def _handle_section_title(self, slide, el, current_top, centered):
        text = get_text(el)
        if not text:
            return current_top
        top = current_top
        if centered:
            top = max(top, 1.2)
        font_size = 36
        shape, tf = _make_text_shape(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, 0.6, auto_size="none"
        )
        _add_paragraph(
            tf,
            text,
            font_size,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT,
        )
        return top + 0.6 + GAP_NORMAL

    def _handle_section_label(self, slide, el, current_top, centered):
        text = get_text(el)
        if not text:
            return current_top
        top = 1.0 if centered else current_top
        shape, tf = _make_text_shape(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, 0.35, auto_size="none"
        )
        _add_paragraph(
            tf,
            text.upper(),
            14,
            bold=True,
            color=self.accent_color,
            alignment=PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT,
        )
        return top + 0.4 + GAP_TIGHT

    def _handle_headline(self, slide, el, current_top, centered):
        text = get_text(el)
        if not text:
            return current_top
        classes = el.get("class", [])
        is_h1 = el.name == "h1"
        is_big_text = "big-text" in classes

        font_size = 48 if is_h1 else 40
        color = MS_CYAN if is_big_text else WHITE
        if is_big_text:
            font_size = 52

        top = current_top
        if centered:
            top = max(top, 1.5)

        height_est = font_size / 72 * 1.5 + 0.1
        shape, tf = _make_text_shape(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, height_est, auto_size="none"
        )
        if centered:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Check for rich text
        rich = get_rich_text(el)
        has_formatting = (
            any(r.get("bold") or r.get("color") for r in rich) if rich else False
        )
        if has_formatting and rich:
            _add_rich_paragraphs(
                tf,
                rich,
                font_size,
                default_color=color,
                alignment=PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT,
                bold=True,
            )
        else:
            _add_paragraph(
                tf,
                text,
                font_size,
                bold=True,
                color=color,
                alignment=PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT,
            )

        gap = 0.50 if centered else GAP_SECTION
        return top + height_est + gap

    def _handle_medium_headline(self, slide, el, current_top, centered):
        text = get_text(el)
        if not text:
            return current_top
        font_size = 36
        height_est = 0.6
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height_est,
            auto_size="none",
        )
        color = WHITE
        style_color = self._resolve_style_color(el)
        if style_color:
            color = style_color
        _add_paragraph(
            tf,
            text,
            font_size,
            bold=True,
            color=color,
            alignment=PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT,
        )
        return current_top + height_est + GAP_SECTION

    def _handle_subhead(self, slide, el, current_top, centered):
        text = get_text(el)
        if not text:
            return current_top

        rich = get_rich_text(el)
        has_formatting = (
            any(r.get("color") or r.get("italic") or r.get("bold") for r in rich)
            if rich
            else False
        )

        height_est = 0.45
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height_est,
            auto_size="none",
        )
        align = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT
        if has_formatting and rich:
            _add_rich_paragraphs(tf, rich, 24, default_color=GRAY_70, alignment=align)
        else:
            _add_paragraph(tf, text, 24, color=GRAY_70, alignment=align)
        return current_top + height_est + GAP_NORMAL

    def _handle_architecture_diagram(self, slide, el, current_top):
        """Render architecture diagram as code-block style."""
        code_el = el.find("code") or el.find("pre") or el
        runs = self._extract_code_runs(code_el)
        if not runs:
            text = get_text(el)
            if not text:
                return current_top
            runs = [{"text": text, "color": CODE_DEFAULT, "bold": False}]

        full_text = "".join(r["text"] for r in runs)
        num_lines = full_text.count("\n") + 1
        height = min(4.5, max(1.0, num_lines * 0.18 + 0.3))

        # Cap to available slide space
        available = SLIDE_HEIGHT - current_top - 0.15
        if available < 0.8:
            available = 0.8
        if height > available:
            height = available

        # Background
        _make_filled_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height,
            fill_color=CODE_BG,
            border_color=BORDER_GRAY,
        )

        # Text overlay
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT + 0.05,
            current_top + 0.05,
            CONTENT_WIDTH - 0.1,
            height - 0.1,
            auto_size="fit_shape",
        )
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        lines_text = full_text.split("\n")
        font_size = 10 if num_lines > 20 else 11 if num_lines > 12 else 12

        self._render_code_runs_to_tf(tf, runs, lines_text, font_size)
        return current_top + height + GAP_NORMAL

    def _render_code_runs_to_tf(self, tf, runs, lines, font_size):
        """Render syntax-colored code runs into a text frame, line by line."""
        # Build a flat character-to-run mapping
        current_pos = 0
        for line_idx, line in enumerate(lines):
            if line_idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            line_end = current_pos + len(line)
            run_pos = 0
            has_content = False
            for r in runs:
                r_start = run_pos
                r_end = run_pos + len(r["text"])
                overlap_start = max(current_pos, r_start)
                overlap_end = min(line_end, r_end)
                if overlap_start < overlap_end:
                    segment = r["text"][overlap_start - r_start : overlap_end - r_start]
                    if segment:
                        pptx_run = p.add_run()
                        pptx_run.text = segment
                        _set_run_font(
                            pptx_run,
                            name=CODE_FONT,
                            size=font_size,
                            bold=r.get("bold", False),
                            color=r["color"],
                        )
                        has_content = True
                run_pos = r_end
            current_pos = line_end + 1  # +1 for \n

            if not has_content:
                pptx_run = p.add_run()
                pptx_run.text = " "
                _set_run_font(
                    pptx_run, name=CODE_FONT, size=font_size, color=CODE_DEFAULT
                )

    def _handle_comparison_table(self, slide, el, current_top):
        """Render comparison-table as native table."""
        table_el = el.find("table")
        if table_el:
            return self._handle_html_table(slide, table_el, current_top)

        # CSS grid style comparison
        headers = el.find_all(class_="header")
        lefts = el.find_all(class_="left")
        rights = el.find_all(class_="right")

        rows_data = []
        if headers:
            rows_data.append([get_text(h) for h in headers])
        for l_el, r_el in zip(lefts, rights):
            rows_data.append([get_text(l_el), get_text(r_el)])

        if not rows_data:
            return current_top

        _, est_h = self._make_table(
            slide,
            rows_data,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            header=bool(headers),
            font_size=12,
            accent_color=self.accent_color,
        )
        return current_top + est_h + GAP_NORMAL

    def _handle_card_grid(self, slide, container, current_top, handled):
        """Render a card grid as a native PowerPoint table."""
        classes = container.get("class", [])

        children = [
            c
            for c in container.children
            if isinstance(c, Tag) and c.name not in ("style", "script")
        ]
        if not children:
            return current_top

        card_classes = {"card", "module-card", "tool-card", "track-card"}
        items = []  # (type, element) pairs
        for child in children:
            child_classes = set(child.get("class", []))
            if child_classes & card_classes or child.find(
                class_=lambda c: c is not None and c in card_classes
            ):
                items.append(("card", child))
            else:
                items.append(("other", child))
            # Mark child AND all descendants as handled to prevent
            # later handlers (body-text, etc.) from double-rendering
            handled.add(id(child))
            for desc in child.find_all(True):
                handled.add(id(desc))

        if not items:
            return current_top

        cols = self._determine_grid_cols(classes, len(items))

        # Separate code blocks (render separately after the table)
        code_blocks_after = []
        table_items = []
        for item_type, item_el in items:
            item_classes = item_el.get("class", [])
            if "code-block" in item_classes:
                code_blocks_after.append(item_el)
            else:
                table_items.append((item_type, item_el))

        if not table_items and not code_blocks_after:
            return current_top

        # Build table for card items
        if table_items:
            actual_cols = min(cols, len(table_items))
            actual_rows = math.ceil(len(table_items) / actual_cols)
            row_height = Inches(1.5)

            table_shape = slide.shapes.add_table(
                actual_rows,
                actual_cols,
                Inches(CONTENT_LEFT),
                Inches(current_top),
                Inches(CONTENT_WIDTH),
                row_height * actual_rows,
            )
            table = table_shape.table

            col_w = CONTENT_WIDTH / actual_cols
            for i in range(actual_cols):
                table.columns[i].width = Inches(col_w)

            for idx, (item_type, item_el) in enumerate(table_items):
                r = idx // actual_cols
                c = idx % actual_cols
                cell = table.cell(r, c)

                if item_type == "card":
                    card_data = self._extract_card_content(item_el)
                    self._populate_card_cell(
                        cell,
                        card_data,
                        title_size=14 if actual_cols >= 3 else 16,
                        body_size=11 if actual_cols >= 3 else 12,
                        accent_color=self.accent_color,
                    )
                else:
                    # Non-card item: just put text in the cell
                    text = get_text(item_el)
                    cell.text = ""
                    tf = cell.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.1)
                    tf.margin_right = Inches(0.1)
                    tf.margin_top = Inches(0.08)
                    tf.margin_bottom = Inches(0.08)
                    if text:
                        _add_paragraph(tf, text, 11, color=GRAY_70)

                _set_cell_fill(cell, DARK_GRAY)
                _set_cell_border(cell, "333333", 6350)

            # Fill empty cells
            for idx in range(len(table_items), actual_rows * actual_cols):
                r = idx // actual_cols
                c = idx % actual_cols
                cell = table.cell(r, c)
                cell.text = ""
                _set_cell_fill(cell, DARK_GRAY)
                _set_cell_no_border(cell)

            grid_height = actual_rows * 1.8 + (actual_rows - 1) * 0.15
            current_top += grid_height + GAP_NORMAL

        # Render code blocks after the table
        for code_el in code_blocks_after:
            current_top = self._handle_code_block(slide, code_el, current_top)

        return current_top

    def _handle_standalone_card(self, slide, el, current_top):
        """Render a standalone card as a single shape."""
        card_data = self._extract_card_content(el)
        _make_card_frame(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            1.2,
            title=card_data["title"],
            body=card_data["body"],
            body_runs=card_data.get("body_runs"),
            title_color=card_data.get("title_color") or self.accent_color,
            accent_color=self.accent_color,
            border_color=BORDER_GRAY,
        )
        return current_top + 1.5 + GAP_NORMAL

    def _handle_principles(self, slide, el, current_top, handled):
        """Render principles as a two-column layout."""
        principles = (
            el.find_all(class_="principle")
            if el.get("class") and "principles-grid" in el.get("class", [])
            else [el]
        )
        if not principles:
            return current_top

        num = len(principles)
        use_two_cols = num >= 4
        col_width = 4.0 if use_two_cols else CONTENT_WIDTH
        col_offset = 4.4 if use_two_cols else 0
        top = current_top
        col_tops = [top, top]

        for i, p_el in enumerate(principles):
            handled.add(id(p_el))
            col = i % 2 if use_two_cols else 0
            left = CONTENT_LEFT + (col * col_offset)
            ct = col_tops[col]

            num_el = p_el.find(
                class_=lambda c: c and ("principle-number" in c or "principle-num" in c)
            )
            content_el = p_el.find(
                class_=lambda c: (
                    c and ("principle-content" in c or "principle-text" in c)
                )
            )

            number = get_text(num_el) if num_el else ""
            title = ""
            body = ""

            if content_el:
                h3 = content_el.find("h3")
                if h3:
                    title = get_text(h3)
                p_tags = content_el.find_all("p")
                body_parts = [get_text(p) for p in p_tags if get_text(p)]
                body = "\n".join(body_parts)
                if not title:
                    strong = content_el.find("strong")
                    if strong:
                        title = get_text(strong)
            else:
                title = get_text(p_el)

            if number and title:
                title = f"{number}  {title}"

            accent = parse_color_from_class(p_el.get("class", [])) or self.accent_color
            self._render_tenet_shape(slide, left, ct, col_width, title, body, accent)
            col_tops[col] = ct + 1.2 + GAP_TIGHT

        return max(col_tops) + GAP_NORMAL

    def _handle_code_block(self, slide, el, current_top):
        """Render code block with syntax highlighting."""
        code_el = el.find("code") or el.find("pre") or el
        runs = self._extract_code_runs(code_el)

        if not runs:
            text = get_text(el)
            if not text:
                return current_top
            runs = [{"text": text, "color": CODE_DEFAULT, "bold": False}]

        full_text = "".join(r["text"] for r in runs)
        lines = full_text.split("\n")
        # Remove trailing empty lines
        while lines and not lines[-1].strip():
            lines.pop()
        num_lines = len(lines)
        height = min(4.5, max(1.0, num_lines * 0.18 + 0.3))

        # Cap to available slide space — never extend past bottom
        available = SLIDE_HEIGHT - current_top - 0.15
        if available < 0.8:
            available = 0.8
        if height > available:
            height = available

        font_size = (
            9
            if num_lines > 25
            else 10
            if num_lines > 18
            else 11
            if num_lines > 12
            else 12
        )

        # Background
        _make_filled_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height,
            fill_color=CODE_BG,
            border_color=BORDER_GRAY,
        )

        # Text — TEXT_TO_FIT_SHAPE lets PowerPoint shrink font if needed
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT + 0.05,
            current_top + 0.05,
            CONTENT_WIDTH - 0.1,
            height - 0.1,
            auto_size="fit_shape",
        )
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        self._render_code_runs_to_tf(tf, runs, lines, font_size)
        return current_top + height + GAP_NORMAL

    def _handle_flow_diagram(self, slide, el, current_top):
        """Render flow diagram as horizontal steps with arrows."""
        _step_classes = {"flow-box", "flow-step", "workflow-step", "step"}
        steps = el.find_all(
            class_=lambda c: c is not None and c in _step_classes
        )
        if not steps:
            return current_top

        num_steps = len(steps)
        max_per_row = 4
        rows = math.ceil(num_steps / max_per_row)
        cols_per_row = min(num_steps, max_per_row)
        box_gap = 0.15
        arrow_width = 0.3
        total_arrow_width = (cols_per_row - 1) * arrow_width if cols_per_row > 1 else 0
        total_gap_width = (cols_per_row - 1) * box_gap if cols_per_row > 1 else 0
        box_width = (CONTENT_WIDTH - total_arrow_width - total_gap_width) / cols_per_row
        box_height = 1.0

        for row_idx in range(rows):
            row_top = current_top + row_idx * (box_height + 0.3)
            start_idx = row_idx * max_per_row
            row_steps = steps[start_idx : start_idx + max_per_row]

            for i, step_el in enumerate(row_steps):
                left = CONTENT_LEFT + i * (box_width + arrow_width + box_gap)

                _sn = {"step-number", "flow-step-number", "workflow-step-number"}
                num_el = (
                    step_el.find(class_=lambda c: c is not None and c in _sn)
                    if step_el
                    else None
                )
                _st = {"flow-step-title", "workflow-step-title", "step-title"}
                title_el = step_el.find(
                    class_=lambda c: c is not None and c in _st
                )
                _sd = {"flow-step-desc", "workflow-step-desc", "step-desc"}
                desc_el = step_el.find(
                    class_=lambda c: c is not None and c in _sd
                )

                number = get_text(num_el) if num_el else ""
                title = get_text(title_el) if title_el else ""
                desc = get_text(desc_el) if desc_el else ""

                if not title and not desc:
                    title = get_text(step_el)
                if number and title:
                    title = f"{number}. {title}"

                _make_card_frame(
                    slide,
                    left,
                    row_top,
                    box_width,
                    box_height,
                    title=title,
                    body=desc,
                    title_size=13,
                    body_size=10,
                    title_color=WHITE,
                    accent_color=self.accent_color,
                    fill_color=DARK_GRAY,
                    border_color=self.accent_color,
                )

                # Arrow between steps
                if i < len(row_steps) - 1:
                    arrow_left = left + box_width + box_gap * 0.3
                    arrow_top = row_top + box_height / 2 - 0.05
                    arrow_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(arrow_left),
                        Inches(arrow_top),
                        Inches(arrow_width * 0.7),
                        Inches(0.1),
                    )
                    arrow_shape.fill.solid()
                    arrow_shape.fill.fore_color.rgb = self.accent_color
                    arrow_shape.line.fill.background()

        total_height = rows * (box_height + 0.3)
        return current_top + total_height + GAP_NORMAL

    def _render_tenet_shape(self, slide, left, top, width, title, body, accent_color):
        """Render a single tenet: accent bar on left + text frame."""
        bar_width = 0.04
        _make_filled_shape(
            slide,
            left,
            top,
            bar_width,
            1.0,
            fill_color=accent_color,
            shape_type=MSO_SHAPE.RECTANGLE,
        )
        shape, tf = _make_text_shape(
            slide,
            left + bar_width + 0.08,
            top,
            width - bar_width - 0.12,
            1.0,
            auto_size="fit_shape",
        )
        if title:
            _add_paragraph(tf, title, 14, bold=True, color=WHITE, space_after=4)
        if body:
            _add_paragraph(tf, body, 11, color=GRAY_70, space_before=2)

    def _handle_tenets(self, slide, elements, current_top, handled):
        """Render tenet elements."""
        if not elements:
            return current_top

        num = len(elements)
        use_two_cols = num >= 4
        col_width = 4.0 if use_two_cols else CONTENT_WIDTH
        col_offset = 4.4 if use_two_cols else 0
        col_tops = [current_top, current_top]

        for i, el in enumerate(elements):
            handled.add(id(el))
            col = i % 2 if use_two_cols else 0
            left = CONTENT_LEFT + (col * col_offset)
            ct = col_tops[col]

            title_el = el.find(class_="tenet-title")
            text_el = el.find(class_="tenet-text")
            title = get_text(title_el) if title_el else ""
            body = get_text(text_el) if text_el else ""
            if not title:
                title = get_text(el)

            accent = parse_color_from_class(el.get("class", [])) or self.accent_color
            self._render_tenet_shape(slide, left, ct, col_width, title, body, accent)
            col_tops[col] = ct + 1.2 + GAP_TIGHT

        return max(col_tops) + GAP_NORMAL

    def _handle_versus(self, slide, el, current_top):
        """Render versus section as two-column cards with 'vs' divider."""
        sides = el.find_all(class_="versus-side")
        if len(sides) < 2:
            return current_top

        left_side = sides[0]
        right_side = sides[1]

        def extract_side(side_el):
            title_el = side_el.find(class_="versus-title")
            title = get_text(title_el) if title_el else ""
            color = parse_color_from_class(side_el.get("class", [])) or MS_ORANGE
            if title_el:
                color = parse_color_from_class(title_el.get("class", [])) or color
            items = []
            feature_list = side_el.find(class_="feature-list") or side_el.find("ul")
            if feature_list:
                for li in feature_list.find_all("li"):
                    items.append(get_text(li))
            return title, color, items

        l_title, l_color, l_items = extract_side(left_side)
        r_title, r_color, r_items = extract_side(right_side)
        if not l_color or l_color == MS_ORANGE:
            l_color = MS_ORANGE
        if not r_color or r_color == MS_ORANGE:
            r_color = MS_GREEN

        col_width = 3.8
        gap = 0.8

        # Left card
        l_body = "\n".join(l_items) if l_items else ""
        _make_card_frame(
            slide,
            CONTENT_LEFT,
            current_top,
            col_width,
            2.0,
            title=l_title,
            body=l_body,
            title_size=18,
            body_size=12,
            title_color=l_color,
            accent_color=l_color,
            fill_color=DARK_GRAY,
            border_color=l_color,
        )

        # "vs" text
        vs_left = CONTENT_LEFT + col_width + (gap - 0.3) / 2
        shape, tf = _make_text_shape(
            slide, vs_left, current_top + 0.8, 0.5, 0.4, auto_size="none"
        )
        _add_paragraph(
            tf, "vs", 18, bold=True, color=GRAY_50, alignment=PP_ALIGN.CENTER
        )

        # Right card
        r_body = "\n".join(r_items) if r_items else ""
        _make_card_frame(
            slide,
            CONTENT_LEFT + col_width + gap,
            current_top,
            col_width,
            2.0,
            title=r_title,
            body=r_body,
            title_size=18,
            body_size=12,
            title_color=r_color,
            accent_color=r_color,
            fill_color=DARK_GRAY,
            border_color=r_color,
        )

        max_items = max(len(l_items), len(r_items), 3)
        height = 0.6 + max_items * 0.3
        return current_top + max(2.0, height) + GAP_NORMAL

    def _handle_html_table(self, slide, el, current_top):
        """Render a standard HTML table as native PowerPoint table."""
        rows_data = []
        has_header = False
        for tr in el.find_all("tr"):
            ths = tr.find_all("th")
            tds = tr.find_all("td")
            if ths:
                has_header = True
                rows_data.append([get_text(th) for th in ths])
            elif tds:
                rows_data.append([get_text(td) for td in tds])

        if not rows_data:
            return current_top

        _, est_h = self._make_table(
            slide,
            rows_data,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            header=has_header,
            font_size=11,
            accent_color=self.accent_color,
        )
        return current_top + est_h + GAP_NORMAL

    def _handle_feature_list(self, slide, el, current_top):
        """Render feature list as single text frame with bullet paragraphs."""
        items = el.find_all("li")
        if not items:
            text = get_text(el)
            if text:
                lines = text.split("\n")
                items_text = [(ln.strip(), None) for ln in lines if ln.strip()]
            else:
                return current_top
        else:
            items_text = []
            for li in items:
                t = get_text(li)
                color = None
                if "\u2713" in t or "✓" in t:
                    color = MS_GREEN
                elif "\u2717" in t or "✗" in t:
                    color = MS_RED
                items_text.append((t, color))

        num_items = len(items_text)
        height = max(0.5, num_items * 0.35)
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height,
            auto_size="fit_shape",
        )

        for text, color in items_text:
            _add_paragraph(
                tf,
                text,
                14,
                color=color or WHITE,
                space_before=2,
                space_after=2,
            )

        return current_top + height + GAP_NORMAL

    def _handle_notification_stack(self, slide, el, current_top):
        """Render notification stack as individual card shapes."""
        notifications = el.find_all(class_="notification")
        if not notifications:
            return current_top

        notif_width = 6.0
        notif_left = (SLIDE_WIDTH - notif_width) / 2
        notif_height = 0.5

        for notif_el in notifications:
            classes = notif_el.get("class", [])
            is_allowed = "allowed" in classes
            is_blocked = "blocked" in classes

            if is_allowed:
                border_color = MS_GREEN
                icon = "\u2713"
                fill = RGBColor(0x0A, 0x1A, 0x0A)
            elif is_blocked:
                border_color = MS_RED
                icon = "\u2717"
                fill = RGBColor(0x1A, 0x0A, 0x0A)
            else:
                border_color = BORDER_GRAY
                icon = "\u2022"
                fill = DARK_GRAY

            title_el = notif_el.find(class_="notification-title")
            body_el = notif_el.find(class_="notification-body")
            title = get_text(title_el) if title_el else ""
            body = get_text(body_el) if body_el else ""

            display_title = f"{icon}  {title}" if title else icon

            _make_card_frame(
                slide,
                notif_left,
                current_top,
                notif_width,
                notif_height,
                title=display_title,
                body=body,
                title_size=12,
                body_size=10,
                title_color=border_color,
                body_color=GRAY_70,
                fill_color=fill,
                border_color=border_color,
            )
            current_top += notif_height + GAP_TIGHT

        return current_top + GAP_NORMAL

    def _handle_stats(self, slide, el, current_top):
        """Render stat-grid, stat-row, or velocity-grid."""
        _stat_exact = {"stat", "velocity-stat"}
        stats = el.find_all(
            class_=lambda c: c is not None and c in _stat_exact
        )
        if not stats:
            return current_top

        num = len(stats)
        col_width = CONTENT_WIDTH / num if num > 0 else CONTENT_WIDTH

        for i, stat_el in enumerate(stats):
            left = CONTENT_LEFT + i * col_width

            _num_classes = {"stat-number", "stat-value", "velocity-number"}
            num_el = stat_el.find(
                class_=lambda c: c is not None and c in _num_classes
            )
            _lbl_classes = {"stat-label", "velocity-label"}
            label_el = stat_el.find(
                class_=lambda c: c is not None and c in _lbl_classes
            )

            number = get_text(num_el) if num_el else ""
            label = get_text(label_el) if label_el else ""

            if not number and not label:
                number = get_text(stat_el)

            shape, tf = _make_text_shape(
                slide,
                left,
                current_top,
                col_width,
                0.9,
                auto_size="none",
            )
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            if number:
                stat_color = parse_color_from_class(stat_el.get("class", [])) or MS_CYAN
                _add_paragraph(
                    tf,
                    number,
                    40,
                    bold=True,
                    color=stat_color,
                    alignment=PP_ALIGN.CENTER,
                    space_after=2,
                )
            if label:
                _add_paragraph(
                    tf,
                    label,
                    12,
                    color=GRAY_70,
                    alignment=PP_ALIGN.CENTER,
                    space_before=2,
                )

        return current_top + 1.0 + GAP_NORMAL

    def _handle_big_stat(self, slide, elements, current_top, handled):
        """Render big-stat elements side by side."""
        if not elements:
            return current_top

        num = len(elements)
        col_width = CONTENT_WIDTH / num if num > 0 else CONTENT_WIDTH

        for i, el in enumerate(elements):
            handled.add(id(el))
            left = CONTENT_LEFT + i * col_width

            num_el = el.find(class_="big-stat-number")
            unit_el = el.find(class_="big-stat-unit")
            label_el = el.find(class_="big-stat-label")

            number = get_text(num_el) if num_el else ""
            unit = get_text(unit_el) if unit_el else ""
            label = get_text(label_el) if label_el else ""

            if not number:
                number = get_text(el)

            shape, tf = _make_text_shape(
                slide,
                left,
                current_top,
                col_width,
                1.2,
                auto_size="none",
            )
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            display_num = f"{number} {unit}".strip() if unit else number
            if display_num:
                _add_paragraph(
                    tf,
                    display_num,
                    56,
                    bold=True,
                    color=MS_CYAN,
                    alignment=PP_ALIGN.CENTER,
                    space_after=4,
                )
            if label:
                _add_paragraph(
                    tf,
                    label,
                    18,
                    color=GRAY_70,
                    alignment=PP_ALIGN.CENTER,
                    space_before=2,
                )

        return current_top + 1.4 + GAP_NORMAL

    def _handle_tier_stack(self, slide, el, current_top):
        """Render tier-stack as sequential tier cards."""
        tiers = el.find_all(class_="tier")
        if not tiers:
            return current_top

        tier_height = 0.6
        for tier_el in tiers:
            label_el = tier_el.find(class_="tier-label")
            title_el = tier_el.find(class_="tier-title")
            desc_el = tier_el.find(class_="tier-desc")
            tokens_el = tier_el.find(class_="tier-tokens")

            label = get_text(label_el) if label_el else ""
            title = get_text(title_el) if title_el else ""
            desc = get_text(desc_el) if desc_el else ""
            tokens = get_text(tokens_el) if tokens_el else ""

            display_title = f"{label}  {title}".strip() if label else title
            body_parts = [p for p in [desc, tokens] if p]
            body = "\n".join(body_parts)

            accent = (
                parse_color_from_class(tier_el.get("class", [])) or self.accent_color
            )

            _make_card_frame(
                slide,
                CONTENT_LEFT,
                current_top,
                CONTENT_WIDTH,
                tier_height,
                title=display_title,
                body=body,
                title_size=14,
                body_size=11,
                accent_color=accent,
                fill_color=DARK_GRAY,
                border_color=accent,
            )
            current_top += tier_height + GAP_TIGHT

        return current_top + GAP_NORMAL

    def _handle_tier_row(self, slide, el, current_top):
        """Render tier-row as native table."""
        rows = (
            el.find_all(class_="tier-row")
            if el.get("class") and "tier-row" not in el.get("class", [])
            else [el]
        )
        rows_data = []

        for row_el in rows:
            name_el = row_el.find(class_="tier-name")
            uses_el = row_el.find(class_="tier-uses")
            cost_el = row_el.find(class_="tier-cost")
            name = get_text(name_el) if name_el else ""
            uses = get_text(uses_el) if uses_el else ""
            cost = get_text(cost_el) if cost_el else ""
            if name or uses or cost:
                rows_data.append([name, uses, cost])

        if not rows_data:
            text = get_text(el)
            if text:
                rows_data.append([text, "", ""])

        if rows_data:
            _, est_h = self._make_table(
                slide,
                rows_data,
                CONTENT_LEFT,
                current_top,
                CONTENT_WIDTH,
                col_widths=[1.8, 4.0, 2.6],
                header=False,
                font_size=12,
                accent_color=self.accent_color,
            )
            return current_top + est_h + GAP_NORMAL
        return current_top

    def _handle_diagram(self, slide, el, current_top):
        """Render diagram with boxes and arrows."""
        boxes = el.find_all(class_="diagram-box")
        if not boxes:
            return current_top

        num = len(boxes)
        box_gap = 0.15
        arrow_width = 0.3
        total_arrows = (num - 1) * arrow_width if num > 1 else 0
        total_gaps = (num - 1) * box_gap if num > 1 else 0
        box_width = (CONTENT_WIDTH - total_arrows - total_gaps) / num
        box_height = 1.0

        for i, box_el in enumerate(boxes):
            left = CONTENT_LEFT + i * (box_width + arrow_width + box_gap)

            title_el = box_el.find(class_="diagram-box-title")
            content_el = box_el.find(class_="diagram-box-content")
            title = get_text(title_el) if title_el else ""
            body = get_text(content_el) if content_el else ""

            if not title and not body:
                title = get_text(box_el)

            _make_card_frame(
                slide,
                left,
                current_top,
                box_width,
                box_height,
                title=title,
                body=body,
                title_size=13,
                body_size=10,
                accent_color=self.accent_color,
                fill_color=DARK_GRAY,
                border_color=self.accent_color,
            )

            if i < num - 1:
                arrow_left = left + box_width + box_gap * 0.3
                arrow_top = current_top + box_height / 2 - 0.05
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(arrow_left),
                    Inches(arrow_top),
                    Inches(arrow_width * 0.7),
                    Inches(0.1),
                )
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = self.accent_color
                arrow_shape.line.fill.background()

        return current_top + box_height + 0.3 + GAP_NORMAL

    def _handle_before_after(self, slide, el, current_top):
        """Render before/after as two-column cards."""
        before_el = el.find(class_="before-card")
        after_el = el.find(class_="after-card")

        col_width = 4.0
        gap = 0.4
        card_height = 2.0

        def extract_comparison(card_el, default_color):
            if not card_el:
                return "", ""
            label_el = card_el.find(class_="comparison-label")
            value_el = card_el.find(class_="comparison-value")
            label = get_text(label_el) if label_el else ""
            value = get_text(value_el) if value_el else ""
            remaining = get_text(card_el)
            for part in [label, value]:
                remaining = remaining.replace(part, "", 1).strip()
            title = label or get_text(card_el.find(["h3", "h4"]) or card_el)
            body_parts = [p for p in [value, remaining] if p]
            return title, "\n".join(body_parts)

        b_title, b_body = extract_comparison(before_el, MS_ORANGE)
        a_title, a_body = extract_comparison(after_el, MS_GREEN)

        _make_card_frame(
            slide,
            CONTENT_LEFT,
            current_top,
            col_width,
            card_height,
            title=b_title or "Before",
            body=b_body,
            title_size=16,
            body_size=12,
            title_color=MS_ORANGE,
            accent_color=MS_ORANGE,
            fill_color=DARK_GRAY,
            border_color=MS_ORANGE,
        )

        _make_card_frame(
            slide,
            CONTENT_LEFT + col_width + gap,
            current_top,
            col_width,
            card_height,
            title=a_title or "After",
            body=a_body,
            title_size=16,
            body_size=12,
            title_color=MS_GREEN,
            accent_color=MS_GREEN,
            fill_color=DARK_GRAY,
            border_color=MS_GREEN,
        )

        return current_top + card_height + GAP_NORMAL

    def _handle_token_display(self, slide, el, current_top):
        """Render token display as monospace text."""
        text = get_text(el)
        if not text:
            return current_top
        lines = text.split("\n")
        num_lines = len(lines)
        height = max(0.5, num_lines * 0.22 + 0.2)

        _make_filled_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height,
            fill_color=CODE_BG,
            border_color=BORDER_GRAY,
        )
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT + 0.1,
            current_top + 0.05,
            CONTENT_WIDTH - 0.2,
            height - 0.1,
            auto_size="fit_shape",
        )
        for line in lines:
            _add_paragraph(tf, line, 11, color=CODE_DEFAULT, font_name=CODE_FONT)

        return current_top + height + GAP_NORMAL

    def _handle_good_bad_pattern(self, slide, good_el, bad_el, current_top):
        """Render good/bad pattern as two columns."""
        col_width = 4.0
        gap = 0.4

        def extract_items(el):
            if not el:
                return []
            items = el.find_all("li")
            if items:
                return [get_text(li) for li in items]
            text = get_text(el)
            return [ln.strip() for ln in text.split("\n") if ln.strip()]

        bad_items = extract_items(bad_el)
        good_items = extract_items(good_el)

        bad_lines = [
            f"\u2717  {item}" if not item.startswith(("\u2717", "\u2718")) else item
            for item in bad_items
        ]
        good_lines = [
            f"\u2713  {item}" if not item.startswith(("\u2713", "\u2714")) else item
            for item in good_items
        ]

        max_items = max(len(bad_lines), len(good_lines), 2)
        height = max(1.0, max_items * 0.3 + 0.4)

        if bad_lines:
            _make_card_frame(
                slide,
                CONTENT_LEFT,
                current_top,
                col_width,
                height,
                title="\u2717  Don't",
                body="\n".join(bad_lines),
                title_size=16,
                body_size=12,
                title_color=MS_RED,
                accent_color=MS_RED,
                fill_color=DARK_GRAY,
                border_color=MS_RED,
            )

        if good_lines:
            _make_card_frame(
                slide,
                CONTENT_LEFT + col_width + gap,
                current_top,
                col_width,
                height,
                title="\u2713  Do",
                body="\n".join(good_lines),
                title_size=16,
                body_size=12,
                title_color=MS_GREEN,
                accent_color=MS_GREEN,
                fill_color=DARK_GRAY,
                border_color=MS_GREEN,
            )

        return current_top + height + GAP_NORMAL

    def _handle_summary_row(self, slide, elements, current_top, handled):
        """Render summary rows as native table."""
        rows_data = []
        for el in elements:
            handled.add(id(el))
            cells = el.find_all(class_="summary-cell")
            if cells:
                rows_data.append([get_text(c) for c in cells])
            else:
                text = get_text(el)
                if text:
                    rows_data.append([text])

        if not rows_data:
            return current_top

        _, est_h = self._make_table(
            slide,
            rows_data,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            header=True,
            font_size=12,
            accent_color=self.accent_color,
        )
        return current_top + est_h + GAP_NORMAL

    def _handle_body_text(self, slide, el, current_top, centered):
        """Render body text paragraph."""
        rich = get_rich_text(el)
        text = get_text(el)
        if not text:
            return current_top

        num_lines = text.count("\n") + 1
        height = max(0.35, num_lines * 0.25 + 0.1)
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height,
            auto_size="fit_shape",
        )

        has_formatting = (
            any(r.get("bold") or r.get("italic") or r.get("color") for r in rich)
            if rich
            else False
        )
        align = PP_ALIGN.CENTER
        if has_formatting and rich:
            _add_rich_paragraphs(tf, rich, 16, default_color=GRAY_70, alignment=align)
        else:
            _add_paragraph(tf, text, 16, color=GRAY_70, alignment=align)

        return current_top + height + GAP_NORMAL

    def _handle_title_meta(self, slide, el, current_top, centered):
        """Render title-meta small text near bottom."""
        text = get_text(el)
        if not text:
            return current_top
        top = max(current_top, 4.5)
        shape, tf = _make_text_shape(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, 0.3, auto_size="none"
        )
        align = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT
        _add_paragraph(tf, text, 12, color=GRAY_50, alignment=align)
        return top + 0.35

    def _handle_highlight_box(self, slide, el, current_top):
        """Render highlight box as filled shape with text."""
        title_el = el.find(["h3", "h4", "strong"])
        title = get_text(title_el) if title_el else ""

        rich = get_rich_text(el)
        text = get_text(el)
        body = text
        if title and body.startswith(title):
            body = body[len(title) :].strip()

        top = min(current_top, SLIDE_HEIGHT - 0.85)

        border_color = self._resolve_style_color(el) or self.accent_color
        classes = el.get("class", [])
        class_color = parse_color_from_class(classes)
        if class_color:
            border_color = class_color

        height = 0.7
        shape = _make_filled_shape(
            slide,
            CONTENT_LEFT,
            top,
            CONTENT_WIDTH,
            height,
            fill_color=DARK_GRAY,
            border_color=border_color,
        )
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        if title:
            _add_paragraph(tf, title, 14, bold=True, color=border_color, space_after=4)
        if body:
            has_formatting = (
                any(r.get("bold") or r.get("italic") or r.get("color") for r in rich)
                if rich
                else False
            )
            if has_formatting and rich:
                body_runs = []
                skipped_title = False
                for r in rich:
                    if not skipped_title and r["text"].strip() == title.strip():
                        skipped_title = True
                        continue
                    body_runs.append(r)
                if body_runs:
                    _add_rich_paragraphs(
                        tf,
                        body_runs,
                        12,
                        default_color=GRAY_70,
                        space_before=2,
                    )
            else:
                _add_paragraph(tf, body, 12, color=GRAY_70, space_before=2)

        return top + 0.8 + GAP_NORMAL

    def _handle_quote(self, slide, el, current_top, centered):
        """Render quote as italic text with attribution."""
        quote_text = ""
        for child in el.children:
            if isinstance(child, Tag):
                child_classes = child.get("class", [])
                if any(x in child_classes for x in ["quote-attribution", "quote-attr"]):
                    continue
                quote_text += get_text(child) + " "
            elif isinstance(child, NavigableString):
                t = str(child).strip()
                if t:
                    quote_text += t + " "
        quote_text = quote_text.strip()
        if not quote_text:
            quote_text = get_text(el)

        attr_el = el.find(
            class_=lambda c: (
                c is not None and c in {"quote-attribution", "quote-attr"}
            )
        )
        attribution = get_text(attr_el) if attr_el else ""
        if attribution and quote_text.endswith(attribution):
            quote_text = quote_text[: -len(attribution)].strip()

        height = 1.0
        shape, tf = _make_text_shape(
            slide,
            CONTENT_LEFT,
            current_top,
            CONTENT_WIDTH,
            height,
            auto_size="fit_shape",
        )

        display = f'"{quote_text}"' if not quote_text.startswith('"') else quote_text
        _add_paragraph(
            tf,
            display,
            24,
            italic=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            space_after=6,
        )
        if attribution:
            attr_display = (
                attribution
                if attribution.startswith("\u2014")
                else f"\u2014 {attribution}"
            )
            _add_paragraph(
                tf,
                attr_display,
                14,
                color=GRAY_50,
                alignment=PP_ALIGN.CENTER,
                space_before=4,
            )

        return current_top + height + GAP_NORMAL

    def _handle_small_text(self, slide, el, current_top, centered):
        """Render small text near bottom of slide."""
        text = get_text(el)
        if not text:
            return current_top
        top = max(current_top, 4.8)
        shape, tf = _make_text_shape(
            slide, CONTENT_LEFT, top, CONTENT_WIDTH, 0.3, auto_size="none"
        )
        align = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT
        _add_paragraph(tf, text, 10, color=GRAY_50, alignment=align)
        return top + 0.3

    # ------------------------------------------------------------------
    # Overflow compression
    # ------------------------------------------------------------------

    def _compress_if_needed(self, slide, slide_num):
        """Proportionally reposition shapes if content overflows."""
        shapes = list(slide.shapes)
        if not shapes:
            return

        actual_bottom = max((s.top + s.height for s in shapes), default=0)
        slide_bottom = Inches(SLIDE_HEIGHT)

        if actual_bottom <= slide_bottom:
            return

        content_start = min(
            (s.top for s in shapes if s.top >= Inches(0.3)),
            default=Inches(0.3),
        )
        content_end = actual_bottom
        content_height = content_end - content_start

        if content_height <= 0:
            return

        usable = slide_bottom - Inches(0.1)
        scale = (usable - content_start) / content_height

        if scale >= 1.0:
            return

        if scale < 0.40:
            scale = 0.40
            self.warnings.append(
                f"Slide {slide_num}: severe overflow \u2014 content scaled to 40%"
            )

        for s in shapes:
            if s.top >= content_start:
                offset = s.top - content_start
                s.top = content_start + int(offset * scale)
                # Scale heights too (not just tops) so shapes don't
                # pile on each other.  Skip tables — their row heights
                # are internal and don't benefit from external scaling.
                if not s.has_table:
                    s.height = int(s.height * scale)

    # ------------------------------------------------------------------
    # Main process_slide
    # ------------------------------------------------------------------

    def process_slide(self, slide_div: Tag, slide_num: int):
        """Process a single slide HTML element."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        set_slide_background(slide)

        centered = self.is_centered(slide_div)
        current_top = 0.6
        handled: set[int] = set()

        # ==========================================================
        # HEADER ZONE: Merge header elements into a single text frame
        # instead of separate absolutely-positioned shapes.
        # This eliminates overlap caused by wrong height estimates.
        # ==========================================================
        header_parts: list[tuple[str, str, dict]] = []

        # Section labels
        for el in slide_div.find_all(class_="section-label"):
            if id(el) in handled:
                continue
            text = get_text(el)
            if text:
                header_parts.append(("label", text, {}))
                handled.add(id(el))

        # Section number (rendered as a label)
        for el in slide_div.find_all(class_="section-number"):
            if id(el) in handled:
                continue
            text = get_text(el)
            if text:
                header_parts.append(("label", text, {}))
                handled.add(id(el))

        # Section title
        for el in slide_div.find_all(class_="section-title"):
            if id(el) in handled:
                continue
            text = get_text(el)
            if text:
                header_parts.append(("section_title", text, {}))
                handled.add(id(el))

        # Headlines (h1, h2.headline, big-text)
        for el in slide_div.find_all(["h1", "h2"]):
            if id(el) in handled:
                continue
            classes = el.get("class", [])
            if "section-title" in classes:
                continue
            if el.name == "h1" or "headline" in classes or "big-text" in classes:
                text = get_text(el)
                if text:
                    is_big = "big-text" in classes
                    is_h1 = el.name == "h1"
                    rich = get_rich_text(el)
                    has_formatting = (
                        any(r.get("bold") or r.get("color") for r in rich)
                        if rich
                        else False
                    )
                    header_parts.append(
                        (
                            "headline",
                            text,
                            {
                                "big": is_big,
                                "h1": is_h1,
                                "rich": rich if has_formatting else None,
                            },
                        )
                    )
                    handled.add(id(el))

        # Medium headline
        for el in slide_div.find_all(class_="medium-headline"):
            if id(el) in handled:
                continue
            text = get_text(el)
            if text:
                style_color = self._resolve_style_color(el)
                header_parts.append(
                    (
                        "medium_headline",
                        text,
                        {"style_color": style_color},
                    )
                )
                handled.add(id(el))

        # Subhead
        for el in slide_div.find_all(class_="subhead"):
            if id(el) in handled:
                continue
            text = get_text(el)
            if text:
                rich = get_rich_text(el)
                has_formatting = (
                    any(
                        r.get("color") or r.get("italic") or r.get("bold") for r in rich
                    )
                    if rich
                    else False
                )
                header_parts.append(
                    (
                        "subhead",
                        text,
                        {"rich": rich if has_formatting else None},
                    )
                )
                handled.add(id(el))

        # Early body-text: include body-text that appears before any
        # content element (grids, tables, code blocks) as an intro
        # paragraph in the header frame.
        if header_parts:
            _content_trigger = {
                "code-block",
                "architecture-diagram",
                "comparison-table",
                "thirds",
                "halves",
                "fourths",
                "grid",
                "grid-2",
                "grid-3",
                "grid-4",
                "grid-5",
                "tools-grid",
                "flow-diagram",
                "workflow",
                "flow",
                "notification-stack",
                "stat-grid",
                "stat-row",
                "velocity-grid",
                "tier-stack",
                "before-after",
                "versus",
                "diagram",
                "principles-grid",
                "card",
                "module-card",
                "tool-card",
            }
            for child in slide_div.children:
                if not isinstance(child, Tag) or child.name in ("style", "script"):
                    continue
                if id(child) in handled:
                    continue
                child_classes = set(child.get("class", []))
                if child_classes & _content_trigger or child.name == "table":
                    break
                if "body-text" in child_classes:
                    bt_text = get_text(child)
                    if bt_text:
                        bt_rich = get_rich_text(child)
                        bt_has_fmt = (
                            any(
                                r.get("bold") or r.get("italic") or r.get("color")
                                for r in bt_rich
                            )
                            if bt_rich
                            else False
                        )
                        header_parts.append(
                            (
                                "body_text",
                                bt_text,
                                {"rich": bt_rich if bt_has_fmt else None},
                            )
                        )
                        handled.add(id(child))

        # ---- Build combined header text frame ----
        if header_parts:
            # Estimate height so we can set a generous bounding box
            _hdr_h = 0.3  # base padding
            for _pt, _tx, _op in header_parts:
                if _pt == "label":
                    _hdr_h += 0.35
                elif _pt == "headline":
                    _lines = max(1, len(_tx) // 35 + 1)
                    _per = 0.75 if _op.get("big") or _op.get("h1") else 0.6
                    _hdr_h += _lines * _per + 0.2
                elif _pt == "section_title":
                    _hdr_h += 0.7
                elif _pt == "medium_headline":
                    _hdr_h += 0.7
                elif _pt == "subhead":
                    _hdr_h += 0.5
                elif _pt == "body_text":
                    _hdr_h += (_tx.count("\n") + 1) * 0.25 + 0.15

            if centered:
                _frame_top = 0.4
                _frame_h = SLIDE_HEIGHT - 0.8
            else:
                _frame_top = 0.5
                _frame_h = _hdr_h

            _hdr_shape, _hdr_tf = _make_text_shape(
                slide,
                CONTENT_LEFT,
                _frame_top,
                CONTENT_WIDTH,
                _frame_h,
                auto_size="fit_shape",
            )
            if centered:
                _hdr_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            for _pt, _tx, _op in header_parts:
                _align = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

                if _pt == "label":
                    _add_paragraph(
                        _hdr_tf,
                        _tx.upper(),
                        14,
                        bold=True,
                        color=self.accent_color,
                        alignment=_align,
                        space_after=4,
                    )
                elif _pt == "headline":
                    _is_big = _op.get("big", False)
                    _is_h1 = _op.get("h1", False)
                    _fsz = 52 if _is_big else 48 if _is_h1 else 40
                    _clr = MS_CYAN if _is_big else WHITE
                    _rich = _op.get("rich")
                    if _rich:
                        _add_rich_paragraphs(
                            _hdr_tf,
                            _rich,
                            _fsz,
                            default_color=_clr,
                            alignment=_align,
                            bold=True,
                            space_before=4,
                            space_after=8,
                        )
                    else:
                        _add_paragraph(
                            _hdr_tf,
                            _tx,
                            _fsz,
                            bold=True,
                            color=_clr,
                            alignment=_align,
                            space_before=4,
                            space_after=8,
                        )
                elif _pt == "section_title":
                    _add_paragraph(
                        _hdr_tf,
                        _tx,
                        36,
                        bold=True,
                        color=WHITE,
                        alignment=_align,
                        space_before=4,
                        space_after=4,
                    )
                elif _pt == "medium_headline":
                    _mh_clr = _op.get("style_color") or WHITE
                    _add_paragraph(
                        _hdr_tf,
                        _tx,
                        32,
                        bold=True,
                        color=_mh_clr,
                        alignment=_align,
                        space_before=4,
                        space_after=8,
                    )
                elif _pt == "subhead":
                    _rich = _op.get("rich")
                    if _rich:
                        _add_rich_paragraphs(
                            _hdr_tf,
                            _rich,
                            20,
                            default_color=GRAY_70,
                            alignment=_align,
                            space_before=4,
                            space_after=4,
                        )
                    else:
                        _add_paragraph(
                            _hdr_tf,
                            _tx,
                            20,
                            color=GRAY_70,
                            alignment=_align,
                            space_before=4,
                            space_after=4,
                        )
                elif _pt == "body_text":
                    _bt_rich = _op.get("rich")
                    if _bt_rich:
                        _add_rich_paragraphs(
                            _hdr_tf,
                            _bt_rich,
                            16,
                            default_color=GRAY_70,
                            alignment=PP_ALIGN.CENTER,
                            space_before=4,
                            space_after=4,
                        )
                    else:
                        _add_paragraph(
                            _hdr_tf,
                            _tx,
                            16,
                            color=GRAY_70,
                            alignment=PP_ALIGN.CENTER,
                            space_before=4,
                            space_after=4,
                        )

            current_top = _frame_top + _frame_h + GAP_SECTION

        # ---- 7. Architecture diagram ----
        for el in slide_div.find_all(class_="architecture-diagram"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_architecture_diagram(slide, el, current_top)

        # ---- 8. Comparison table ----
        for el in slide_div.find_all(class_="comparison-table"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_comparison_table(slide, el, current_top)

        # ---- 9. Card grids ----
        grid_classes = {
            "thirds",
            "halves",
            "fourths",
            "grid",
            "grid-2",
            "grid-3",
            "grid-4",
            "grid-5",
            "tools-grid",
            "split",
        }
        exclude_grids = {"principles-grid", "stat-grid", "velocity-grid"}
        for el in slide_div.find_all(
            class_=lambda c: (
                c is not None and c in grid_classes and c not in exclude_grids
            )
        ):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_card_grid(slide, el, current_top, handled)

        # ---- 10. Standalone cards ----
        card_classes = {"card", "module-card", "tool-card", "track-card"}
        for el in slide_div.find_all(
            class_=lambda c: c is not None and c in card_classes
        ):
            if id(el) in handled:
                continue
            parent = el.parent
            if parent and parent.get("class"):
                parent_classes = set(parent.get("class", []))
                if parent_classes & grid_classes:
                    continue
            handled.add(id(el))
            current_top = self._handle_standalone_card(slide, el, current_top)

        # ---- 11. Principles ----
        for el in slide_div.find_all(class_="principles-grid"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_principles(slide, el, current_top, handled)

        for el in slide_div.find_all(class_="principle"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            title_el = el.find(
                class_=lambda c: (
                    c and "principle" in c and ("content" in c or "text" in c)
                )
            )
            title = get_text(title_el) if title_el else get_text(el)
            self._render_tenet_shape(
                slide,
                CONTENT_LEFT,
                current_top,
                CONTENT_WIDTH,
                title,
                "",
                self.accent_color,
            )
            current_top += 1.2 + GAP_TIGHT

        # ---- 12. Code blocks ----
        for el in slide_div.find_all(class_="code-block"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_code_block(slide, el, current_top)

        # ---- 13. Flow diagrams ----
        for el in slide_div.find_all(
            class_=lambda c: (
                c is not None and c in {"flow-diagram", "workflow", "flow"}
            )
        ):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_flow_diagram(slide, el, current_top)

        # ---- 14. Tenets ----
        tenet_elements = slide_div.find_all(class_="tenet")
        unhandled_tenets = [el for el in tenet_elements if id(el) not in handled]
        if unhandled_tenets:
            current_top = self._handle_tenets(
                slide, unhandled_tenets, current_top, handled
            )

        # ---- 15. Versus ----
        for el in slide_div.find_all(class_="versus"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_versus(slide, el, current_top)

        # ---- 16. Tables ----
        for el in slide_div.find_all("table"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_html_table(slide, el, current_top)

        # ---- 17. Feature lists ----
        for el in slide_div.find_all(class_="feature-list"):
            if id(el) in handled:
                continue
            # Skip if inside a versus
            parent = el.parent
            inside_versus = False
            while parent:
                if parent.get("class") and "versus" in parent.get("class", []):
                    inside_versus = True
                    break
                parent = parent.parent
            handled.add(id(el))
            if not inside_versus:
                current_top = self._handle_feature_list(slide, el, current_top)

        # ---- 18. Notification stacks ----
        for el in slide_div.find_all(class_="notification-stack"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_notification_stack(slide, el, current_top)

        # ---- 19. Stats ----
        for el in slide_div.find_all(
            class_=lambda c: (
                c is not None and c in {"stat-grid", "stat-row", "velocity-grid"}
            )
        ):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_stats(slide, el, current_top)

        # ---- 20. Big stats ----
        big_stats = [
            el for el in slide_div.find_all(class_="big-stat") if id(el) not in handled
        ]
        if big_stats:
            current_top = self._handle_big_stat(slide, big_stats, current_top, handled)

        # ---- 21. Tier stacks ----
        for el in slide_div.find_all(class_="tier-stack"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_tier_stack(slide, el, current_top)

        # ---- 22. Tier rows ----
        tier_rows = [
            el for el in slide_div.find_all(class_="tier-row") if id(el) not in handled
        ]
        if tier_rows:
            rows_data = []
            for el in tier_rows:
                handled.add(id(el))
                name_el = el.find(class_="tier-name")
                uses_el = el.find(class_="tier-uses")
                cost_el = el.find(class_="tier-cost")
                rows_data.append(
                    [
                        get_text(name_el) if name_el else "",
                        get_text(uses_el) if uses_el else "",
                        get_text(cost_el) if cost_el else "",
                    ]
                )
            if rows_data:
                _, est_h = self._make_table(
                    slide,
                    rows_data,
                    CONTENT_LEFT,
                    current_top,
                    CONTENT_WIDTH,
                    col_widths=[1.8, 4.0, 2.6],
                    header=False,
                    font_size=12,
                    accent_color=self.accent_color,
                )
                current_top += est_h + GAP_NORMAL

        # ---- 23. Diagrams ----
        for el in slide_div.find_all(class_="diagram"):
            if id(el) in handled:
                continue
            classes = el.get("class", [])
            if "architecture-diagram" in classes or "flow-diagram" in classes:
                continue
            handled.add(id(el))
            current_top = self._handle_diagram(slide, el, current_top)

        # ---- 24. Before/After ----
        for el in slide_div.find_all(class_="before-after"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_before_after(slide, el, current_top)

        # ---- 25. Token display ----
        for el in slide_div.find_all(class_="token-display"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_token_display(slide, el, current_top)

        # ---- 26. Good/bad patterns ----
        good_el = None
        bad_el = None
        for el in slide_div.find_all(class_="good-pattern"):
            if id(el) not in handled:
                good_el = el
                handled.add(id(el))
                break
        for el in slide_div.find_all(class_="bad-pattern"):
            if id(el) not in handled:
                bad_el = el
                handled.add(id(el))
                break
        if good_el or bad_el:
            current_top = self._handle_good_bad_pattern(
                slide, good_el, bad_el, current_top
            )

        # ---- 27. Summary rows ----
        summary_rows = [
            el
            for el in slide_div.find_all(class_="summary-row")
            if id(el) not in handled
        ]
        if summary_rows:
            current_top = self._handle_summary_row(
                slide, summary_rows, current_top, handled
            )

        # ---- 28. Body text ----
        for el in slide_div.find_all(class_="body-text"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_body_text(slide, el, current_top, centered)

        # ---- 29. Title meta ----
        for el in slide_div.find_all(class_="title-meta"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_title_meta(slide, el, current_top, centered)

        # ---- 30. Highlight boxes ----
        for el in slide_div.find_all(class_="highlight-box"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_highlight_box(slide, el, current_top)

        # ---- 31. Quotes ----
        for el in slide_div.find_all(class_="quote"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_quote(slide, el, current_top, centered)

        # ---- 32. Small text ----
        for el in slide_div.find_all(class_="small-text"):
            if id(el) in handled:
                continue
            handled.add(id(el))
            current_top = self._handle_small_text(slide, el, current_top, centered)

        # ---- 33. Fallback: unhandled elements with text ----
        for el in slide_div.find_all(True, recursive=False):
            if id(el) in handled:
                continue
            if el.name in ("style", "script"):
                continue
            text = get_text(el)
            if not text or len(text) < 3:
                continue
            # Check if any children were handled
            has_handled_children = any(
                id(child) in handled for child in el.find_all(True)
            )
            if has_handled_children:
                continue

            handled.add(id(el))
            shape, tf = _make_text_shape(
                slide,
                CONTENT_LEFT,
                current_top,
                CONTENT_WIDTH,
                0.4,
                auto_size="fit_shape",
            )
            _add_paragraph(tf, text, 14, color=GRAY_50)
            current_top += 0.5 + GAP_NORMAL

        # ---- Overflow compression ----
        self._compress_if_needed(slide, slide_num)

    # ------------------------------------------------------------------
    # Convert and save
    # ------------------------------------------------------------------

    def convert(self) -> Presentation:
        """Convert the HTML to a PowerPoint presentation."""
        slides = self.extract_slides()
        if not slides:
            warnings.warn(
                "No slides found in HTML. Check for <div class='slide'> "
                "or <section class='slide'> elements."
            )
        for i, slide_div in enumerate(slides):
            self.process_slide(slide_div, i + 1)
        return self.prs

    def save(self, output_path: str):
        """Save the presentation to a file."""
        self.prs.save(output_path)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main():
    parser = argparse.ArgumentParser(
        description="Convert HTML decks to PowerPoint (v2 - semantic layout)"
    )
    parser.add_argument("input", help="Input HTML file")
    parser.add_argument("output", nargs="?", help="Output PPTX file")
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output) if args.output else input_path.with_suffix(".pptx")

    html = input_path.read_text(encoding="utf-8")
    print(f"Converting: {input_path}")

    converter = HTMLToPPTXConverterV2(html)
    converter.convert()
    converter.save(str(output_path))

    slide_count = len(converter.prs.slides)
    print(f"Output: {output_path}")
    print(f"Done! Created {output_path} ({slide_count} slides)")

    if converter.warnings:
        print(f"\nWarnings ({len(converter.warnings)}):", file=sys.stderr)
        for w in converter.warnings:
            print(f"  {w}", file=sys.stderr)


if __name__ == "__main__":
    main()
